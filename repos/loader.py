from langchain_community.document_loaders import (
    WebBaseLoader,
    PyPDFLoader,
    Docx2txtLoader,
)
from langchain_community.document_loaders.csv_loader import CSVLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from config import vector_store

import os
import ssl
import urllib3
import requests
from dotenv import load_dotenv

load_dotenv()

# SSL configuration
ssl._create_default_https_context = ssl._create_unverified_context
os.environ["CURL_CA_BUNDLE"] = ""
os.environ["REQUESTS_CA_BUNDLE"] = ""

urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
requests.packages.urllib3.disable_warnings()

session = requests.Session()
session.verify = False


def load_pptx(file_path):
    presentation = Presentation(file_path)
    documents = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())

        if text_parts:
            documents.append(
                Document(
                    page_content="\n".join(text_parts),
                    metadata={
                        "source": file_path,
                        "slide": slide_number,
                    },
                )
            )

    return documents


def load_and_index_documents(source, source_type):

    if source_type == "web":
        loader = WebBaseLoader(
            web_paths=(source,),
            requests_kwargs={"verify": False},
        )

    elif source_type == "pdf":
        loader = PyPDFLoader(source)

    elif source_type == "docx":
        loader = Docx2txtLoader(source)

    elif source_type == "csv":
        loader = CSVLoader(source)

    elif source_type == "pptx":
        docs = load_pptx(source)

    else:
        raise ValueError(f"Unsupported source_type: {source_type}")

    # PPTX is already loaded above
    if source_type != "pptx":
        docs = loader.load()

    print(f"Loaded {len(docs)} documents from {source_type}")

    if docs:
        print("First document content preview:")
        print(f"Page content: {docs[0].page_content}...")
        print(f"Document metadata: {docs[0].metadata}")

    # Split documents into chunks
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
    )

    all_splits = text_splitter.split_documents(docs)

    print("\n=== RecursiveCharacterTextSplitter Results ===")
    print(f"Split {len(docs)} documents into {len(all_splits)} chunks")

    for i, chunk in enumerate(all_splits[:3]):
        print(f"\nChunk {i + 1}:")
        print(f"Content: {chunk.page_content[:200]}...")
        print(f"Length: {len(chunk.page_content)} characters")

    # Store embeddings in vector database
    vector_store.add_documents(documents=all_splits)

    return len(all_splits)